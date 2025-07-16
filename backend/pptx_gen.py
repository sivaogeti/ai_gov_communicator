from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os

def generate_rich_pptx(prompt, image_path=None, script=None, audio_path=None, output_dir="outputs"):
    prs = Presentation()

    # Slide 1: Title Slide
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_title.shapes.title
    subtitle = slide_title.placeholders[1]
    title.text = "AI-Powered Communication"
    subtitle.text = prompt

    # Slide 2: Image
    if image_path and os.path.exists(image_path):
        slide_image = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        left = Inches(1)
        top = Inches(1)
        height = Inches(5.5)
        slide_image.shapes.add_picture(image_path, left, top, height=height)
        title_shape = slide_image.shapes.title
        if title_shape:
            title_shape.text = "Generated Image"

    # Slide 3: Script as bullet points
    if script:
        slide_script = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        slide_script.shapes.title.text = "Narration Script"
        content = slide_script.placeholders[1]
        for line in script.splitlines():
            if line.strip():
                content.text += f"\n• {line.strip()}"

    # Slide 4: Optional Audio path info
    if audio_path and os.path.exists(audio_path):
        slide_audio = prs.slides.add_slide(prs.slide_layouts[1])
        slide_audio.shapes.title.text = "Audio Narration"
        body = slide_audio.placeholders[1]
        body.text = f"The audio narration has been generated and saved to:\n{audio_path}\n\n[Currently not embedded in PPT due to library limitations]"

    os.makedirs(output_dir, exist_ok=True)
    pptx_path = os.path.join(output_dir, "generated_rich_content.pptx")
    prs.save(pptx_path)
    return pptx_path
